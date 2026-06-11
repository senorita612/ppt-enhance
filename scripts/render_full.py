"""把 15 页大纲渲染成完整原生 PPTX。"""
import json
from pathlib import Path

from pptx import Presentation
from ppt_enhance.schemas.outline import SlideOutline
from ppt_enhance.builder.layout_engine import render_outline_to_slide

out_dir = Path(".ppt_enhance_cache/real_notebooklm/outlines")
files = sorted(out_dir.glob("outline_p*.json"))

prs = Presentation()
if prs.slides:
    rid = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rid)
    del prs.slides._sldIdLst[0]

for f in files:
    ol = SlideOutline.model_validate_json(f.read_text())
    render_outline_to_slide(prs, ol)

out = Path("data/samples/real_notebooklm_output/real_regen_full.pptx")
prs.save(str(out))
print("saved", out, "slides:", len(files))
