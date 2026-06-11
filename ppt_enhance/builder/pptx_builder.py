"""SlideIR → PPTX 高保真重建（坐标锚定，只改字不改框）."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ppt_enhance.builder.font_utils import estimate_font_size
from ppt_enhance.schemas.slide_ir import ElementType, SlideIR

# 16:9 默认幻灯片尺寸（英寸）
DEFAULT_SLIDE_W = 13.333
DEFAULT_SLIDE_H = 7.5


def _emu(value_inches: float) -> Emu:
    return Inches(value_inches)


def _map_bbox_to_slide(
    bbox,
    page_w: float,
    page_h: float,
    slide_w: float = DEFAULT_SLIDE_W,
    slide_h: float = DEFAULT_SLIDE_H,
) -> tuple[float, float, float, float]:
    """将页面像素 bbox 映射到幻灯片英寸坐标."""
    left = (bbox.x0 / page_w) * slide_w
    top = (bbox.y0 / page_h) * slide_h
    width = (bbox.width / page_w) * slide_w
    height = (bbox.height / page_h) * slide_h
    return left, top, width, height


def _set_slide_size(prs: Presentation, page_w: float, page_h: float) -> tuple[float, float]:
    ratio = page_w / page_h if page_h else 16 / 9
    if ratio > 1.65:
        slide_w, slide_h = 13.333, 7.5
    elif ratio > 1.45:
        slide_w, slide_h = 13.333, 8.333
    elif ratio > 1.0:
        slide_w, slide_h = 11.0, 8.5
    else:
        slide_h = 10.0
        slide_w = slide_h * ratio
    prs.slide_width = _emu(slide_w)
    prs.slide_height = _emu(slide_h)
    return slide_w, slide_h


def _hex_to_rgbcolor(hex_str: str | None, default: tuple[int, int, int]) -> RGBColor:
    if hex_str and hex_str.startswith("#") and len(hex_str) == 7:
        try:
            return RGBColor(int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16))
        except ValueError:
            pass
    return RGBColor(*default)


def _set_slide_background(slide, hex_color: str | None) -> None:
    """设置幻灯片背景为纯色填充。"""
    if not hex_color:
        return
    rgb = _hex_to_rgbcolor(hex_color, (255, 255, 255))
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_pptx(
    slide_ir: SlideIR,
    output_path: str | Path,
    use_background: bool = True,
) -> Path:
    """从 SlideIR 构建 PPTX 文件.

    use_background=True：原图整页铺底 + 透明文字框（高保真但文字是双份的）。
    use_background=False：真·分层重建——纯色背景 + 采样文字色 + 装饰图片对象，
        没有原图像素文字，是真正可编辑的全新 PPTX。
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    # 删除默认空白页
    if prs.slides:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    blank_layout = prs.slide_layouts[6]  # 空白布局

    for page in slide_ir.pages:
        slide_w, slide_h = _set_slide_size(prs, page.width, page.height)
        slide = prs.slides.add_slide(blank_layout)

        # 背景：高保真模式铺原图；分层模式铺「无字背景版」(保留装饰/框线)，回退纯色
        if use_background and page.render_path and Path(page.render_path).exists():
            slide.shapes.add_picture(
                page.render_path,
                _emu(0), _emu(0),
                width=_emu(slide_w),
                height=_emu(slide_h),
            )
        elif not use_background:
            _set_slide_background(slide, page.background_color)
            if page.background_image and Path(page.background_image).exists():
                slide.shapes.add_picture(
                    page.background_image,
                    _emu(0), _emu(0),
                    width=_emu(slide_w),
                    height=_emu(slide_h),
                )

        for elem in page.elements:
            left, top, width, height = _map_bbox_to_slide(
                elem.bbox, page.width, page.height, slide_w, slide_h
            )

            if elem.type == ElementType.IMAGE and elem.image_path and Path(elem.image_path).exists():
                # 高保真模式已有原图铺底；分层模式由无字背景版承载装饰，故均跳过
                if not use_background and elem.metadata.get("layer") != "decoration":
                    slide.shapes.add_picture(
                        elem.image_path,
                        _emu(left), _emu(top),
                        width=_emu(width),
                        height=_emu(height),
                    )
                continue

            if not elem.is_textual or not elem.final_text.strip():
                continue

            txBox = slide.shapes.add_textbox(
                _emu(left), _emu(top), _emu(width), _emu(height)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            p.text = elem.final_text
            p.alignment = PP_ALIGN.LEFT

            # 字号按"映射到画布后的框高度"推算（磅），而非原始像素高度。
            # font_size_hint 是像素量纲，直接当磅值会使字号严重偏大、文字溢出重叠。
            box_h_pt = height * 72.0  # 英寸 → 磅
            char_count = max(1, len(elem.final_text))
            box_w_pt = width * 72.0
            # 单行可容纳字数估算（CJK 约等宽，宽≈高）；多行则按高度均摊
            est_lines = max(1, round(char_count * box_h_pt / max(box_w_pt, 1e-6)))
            font_size = estimate_font_size(box_h_pt, line_count=est_lines, min_pt=8.0)
            # 不超过框高，留 5% 余量防溢出
            font_size = min(font_size, box_h_pt * 0.95)
            p.font.size = Pt(max(8.0, font_size))
            p.font.name = "Microsoft YaHei"
            # 分层模式用采样文字色；高保真模式保持黑色
            if use_background:
                p.font.color.rgb = RGBColor(0, 0, 0)
            else:
                p.font.color.rgb = _hex_to_rgbcolor(elem.text_color, (0, 0, 0))

            # 背景模式下文本框设透明底
            if use_background:
                txBox.fill.background()
                txBox.line.fill.background()

    prs.save(str(output_path))
    return output_path
