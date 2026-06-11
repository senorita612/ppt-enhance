"""布局引擎（坐标锚定版）：SlideOutline → 原生 PPT 形状，按原图真实坐标摆放。

每个元素的位置来自 OCR 真实坐标(px)，映射到画布英寸 → 版面与原图一致。
所有元素均为 PPT 原生对象(圆角矩形/连接符/文本框)，100% 可编辑。
"""

from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SLIDE_W = 13.333
SLIDE_H = 7.5


def _hex(h: str, default=(255, 255, 255)) -> RGBColor:
    if h and h.startswith("#") and len(h) == 7:
        try:
            return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))
        except ValueError:
            pass
    return RGBColor(*default)


def _lighten(h: str, factor: float = 0.16) -> RGBColor:
    c = _hex(h, (20, 40, 60))
    return RGBColor(
        min(255, int(c[0] + (255 - c[0]) * factor)),
        min(255, int(c[1] + (255 - c[1]) * factor)),
        min(255, int(c[2] + (255 - c[2]) * factor)),
    )


class _Mapper:
    """像素 → 画布英寸映射（保持原图宽高比，居中适配）。"""

    def __init__(self, page_w, page_h):
        self.pw = page_w or 1
        self.ph = page_h or 1
        # 等比缩放铺满画布（页面与画布同为16:9时即完全贴合）
        self.sx = SLIDE_W / self.pw
        self.sy = SLIDE_H / self.ph

    def rect(self, bbox_px):
        x0, y0, x1, y1 = bbox_px
        return (x0 * self.sx, y0 * self.sy, (x1 - x0) * self.sx, (y1 - y0) * self.sy)

    def pt_from_h(self, px_h):
        """像素高 → 字号磅（画布上的视觉高度）。"""
        return max(8.0, px_h * self.sy * 72.0)


def _set_alpha(color_elem, opacity: float):
    """给 solidFill 的颜色加透明度（python-pptx 无 API，写 XML）。opacity 0~1。"""
    srgb = color_elem._xFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
    srgb.append(alpha)


def _apply_shadow(shape):
    """加柔和外阴影（写 XML）。"""
    spPr = shape._element.spPr
    effLst = spPr.find(qn("a:effectLst"))
    if effLst is None:
        effLst = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(effLst)
    sh = effLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "63500", "dist": "38100", "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = clr.makeelement(qn("a:alpha"), {"val": "40000"})
    clr.append(alpha)
    sh.append(clr)
    effLst.append(sh)


def _set_radius(shape, ratio: float):
    """设置圆角矩形的圆角大小（adj 值，0~50000）。"""
    try:
        shape.adjustments[0] = max(0.0, min(0.5, ratio))
    except (IndexError, Exception):
        pass


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex(color, (14, 42, 63))


def _draw_grid(slide, accent_hex, *, step=1.1, opacity=0.22):
    """绘制淡网格纹理（蓝图质感）。线条用强调色低透明度，纯原生。"""
    n_v = int(SLIDE_W / step) + 1
    n_h = int(SLIDE_H / step) + 1
    col = _hex(accent_hex, (39, 166, 139))
    for i in range(1, n_v):
        x = i * step
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(x), Inches(0), Inches(x), Inches(SLIDE_H))
        ln.line.color.rgb = col
        ln.line.width = Pt(0.5)
        _line_alpha(ln, opacity)
    for j in range(1, n_h):
        y = j * step
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(0), Inches(y), Inches(SLIDE_W), Inches(y))
        ln.line.color.rgb = col
        ln.line.width = Pt(0.5)
        _line_alpha(ln, opacity)


def _line_alpha(connector, opacity):
    """给线条颜色加透明度。"""
    ln = connector.line._get_or_add_ln()
    srgb = ln.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
    if srgb is None:
        fill = ln.find(qn("a:solidFill"))
        if fill is not None:
            srgb = fill.find(qn("a:srgbClr"))
    if srgb is not None:
        a = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
        srgb.append(a)


def _textbox(slide, text, l, t, w, h, *, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(max(0.2, h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "PingFang SC"
    r.font.color.rgb = _hex(color, (255, 255, 255))
    return tb


def _card(slide, l, t, w, h, style):
    """按 StyleSpec 渲染卡片：solid/translucent/outline + 圆角 + 阴影。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    if style.card_style == "outline":
        # 纯描边也给一层极淡填充 + 阴影，让卡片"浮起来"，避免单薄透明
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex(style.card_fill, (16, 32, 48))
        _set_alpha(shp.fill.fore_color, 0.45)
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex(style.card_fill, (28, 51, 70))
        if style.card_style == "translucent":
            _set_alpha(shp.fill.fore_color, style.card_opacity)
    shp.line.color.rgb = _hex(style.card_border, (39, 166, 139))
    shp.line.width = Pt(style.card_border_width)
    shp.shadow.inherit = False
    _set_radius(shp, style.card_radius)
    # outline 风格强制加阴影提升层次
    if style.card_shadow or style.card_style == "outline":
        _apply_shadow(shp)
    return shp


def _enable_autoshrink(text_frame):
    """开启 PowerPoint 原生 normAutofit：文本溢出时自动缩小字号，兜底防溢出。"""
    try:
        bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
        if bodyPr is None:
            return
        for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
            ex = bodyPr.find(qn(tag))
            if ex is not None:
                bodyPr.remove(ex)
        bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))
    except Exception:
        pass


def _cjk_aware_len(text: str) -> float:
    """文本视觉宽度（CJK算1，ASCII算0.55）。"""
    total = 0.0
    for ch in text:
        total += 1.0 if ord(ch) > 0x2E80 else 0.55
    return total


def _fit_font_size(segments, card_w_in, card_h_in, *, max_pt, min_pt=8.0,
                   margin_in=0.28):
    """二分搜索使所有文字段落恰好填满卡片的最大基准字号(磅)。"""
    avail_w = max(0.5, card_w_in - margin_in)
    avail_h = max(0.3, card_h_in - margin_in)
    avail_w_pt = avail_w * 72.0

    def height_at(base):
        total_lines = 0.0
        for text, ratio in segments:
            if not text:
                continue
            fs = base * ratio
            char_w_pt = fs * 0.92
            chars_per_line = max(1.0, avail_w_pt / char_w_pt)
            lines = max(1.0, _cjk_aware_len(text) / chars_per_line)
            total_lines += lines * ratio
        return total_lines * base * 1.34 / 72.0

    lo, hi = min_pt, max_pt
    best = min_pt
    for _ in range(20):
        base = (lo + hi) / 2
        if height_at(base) <= avail_h:
            best = base
            lo = base
        else:
            hi = base
    return best


# 各层级字号倍率与绝对上限（磅），防止短标题被放飞
_RATIO = {"heading": 1.25, "subtext": 0.8, "body": 1.0}
_CAP = {"heading": 17.0, "subtext": 12.0, "body": 14.0}


def _fill_card_text(shape, node, text_color, accent, style):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.1 if not style.header_band else 0.06)
    tf.margin_bottom = Inches(0.08)
    _enable_autoshrink(tf)  # PowerPoint 原生「溢出自动缩字」兜底

    cw = shape.width / 914400.0
    ch = shape.height / 914400.0
    # body 行：区分公式与文字
    body_lines = node.body_lines or ([type("BL", (), {"text": node.body, "kind": "text"})()]
                                     if node.body else [])
    body_text = " ".join(bl.text for bl in body_lines if bl.kind == "text").strip()
    segs = []
    if node.heading:
        segs.append((node.heading, _RATIO["heading"]))
    if node.subtext:
        segs.append((node.subtext, _RATIO["subtext"]))
    if body_text:
        segs.append((body_text, _RATIO["body"]))
    if not segs and not body_lines:
        return
    base = _fit_font_size(segs, cw, ch, max_pt=15.0, min_pt=6.5) if segs else 12.0

    def sized(kind):
        return min(base * _RATIO[kind], _CAP[kind])

    first = True
    if node.heading:
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = node.heading
        r.font.size = Pt(sized("heading")); r.font.bold = True
        r.font.color.rgb = _hex(accent, (39, 166, 139)); r.font.name = style.font_name
        first = False
    if node.subtext:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        r = p.add_run(); r.text = node.subtext
        r.font.size = Pt(sized("subtext")); r.font.italic = True
        r.font.color.rgb = _hex(text_color, (210, 210, 210)); r.font.name = style.font_name
        first = False
    for bl in body_lines:
        if not bl.text.strip():
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(base * 0.25)
        first = False
        if bl.kind == "formula" and _add_formula_run(p, bl.text, text_color):
            continue  # 成功插入 OMML 公式
        r = p.add_run(); r.text = bl.text
        r.font.size = Pt(sized("body"))
        r.font.color.rgb = _hex(text_color, (230, 230, 230)); r.font.name = style.font_name


def _add_formula_run(paragraph, latex_text, text_color) -> bool:
    """把公式行转成 OMML 并以 PowerPoint 标准结构插入段落。成功 True，失败 False。"""
    try:
        from ppt_enhance.builder.formula_render import (
            latex_to_omml, build_alt_content, clean_inline_latex,
        )
        omml = latex_to_omml(latex_text)
        if not omml:
            return False
        fallback = clean_inline_latex(latex_text)
        paragraph._p.append(build_alt_content(omml, fallback))
        return True
    except Exception:
        return False


def _arrowhead(connector):
    ln = connector.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = ln.makeelement(qn("a:tailEnd"), {}); ln.append(tail)
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")


def _connect(slide, src_rect, dst_rect, style):
    """在两个矩形(l,t,w,h)间连箭头。对齐则直线，错位则用肘形(正交)连接，避免丑陋斜线。"""
    sl, st, sw, sh = src_rect
    dl, dt, dw, dh = dst_rect
    scx, scy = sl + sw / 2, st + sh / 2
    dcx, dcy = dl + dw / 2, dt + dh / 2
    dx, dy = dcx - scx, dcy - scy
    if abs(dx) >= abs(dy):
        if dx >= 0:
            bx, by, ex, ey = sl + sw, scy, dl, dcy
        else:
            bx, by, ex, ey = sl, scy, dl + dw, dcy
        misalign = abs(dy)
    else:
        if dy >= 0:
            bx, by, ex, ey = scx, st + sh, dcx, dt
        else:
            bx, by, ex, ey = scx, st, dcx, dt + dh
        misalign = abs(dx)
    # 错位明显(>0.6英寸)用肘形正交连接，否则直线
    kind = MSO_CONNECTOR.ELBOW if misalign > 0.6 else MSO_CONNECTOR.STRAIGHT
    c = slide.shapes.add_connector(kind, Inches(bx), Inches(by), Inches(ex), Inches(ey))
    c.line.color.rgb = _hex(style.arrow_color, (39, 166, 139))
    c.line.width = Pt(style.arrow_width)
    _arrowhead(c)


def _place_keep_regions(slide, outline, m):
    """把保留的复杂图形区域裁原始像素，作为图片对象放回原位置。"""
    if not outline.keep_regions or not outline.source_image_path:
        return
    import cv2
    if not Path(outline.source_image_path).exists():
        return
    img = cv2.imread(outline.source_image_path)
    if img is None:
        return
    cache_dir = Path(".ppt_enhance_cache/keep_regions")
    cache_dir.mkdir(parents=True, exist_ok=True)
    for i, kr in enumerate(outline.keep_regions):
        x0, y0, x1, y1 = [int(v) for v in kr.bbox_px]
        if x1 <= x0 or y1 <= y0:
            continue
        crop = img[y0:y1, x0:x1]
        if crop.size == 0:
            continue
        out = cache_dir / f"p{outline.page_no:04d}_kr{i}.png"
        cv2.imwrite(str(out), crop)
        l, t, w, h = m.rect(kr.bbox_px)
        slide.shapes.add_picture(str(out), Inches(l), Inches(t), Inches(w), Inches(h))


def _multiline_textbox(slide, lines, l, t, w, h, *, size, color, bold=False,
                       italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """多行文本框：每个元素一行（真正换行，而非空格连接）。"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(max(0.2, h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.12
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = "PingFang SC"; r.font.color.rgb = _hex(color, (255, 255, 255))
    return tb


def _render_cover(prs, slide, outline, m, style):
    """封面专属布局：左侧大标题文字 + 右侧保留插画。竖向堆叠，避免重叠。"""
    _place_keep_regions(slide, outline, m)

    has_illust = bool(outline.keep_regions)
    left = 0.7  # 统一左边距，留白
    text_right = (SLIDE_W * 0.52 if has_illust else SLIDE_W - 0.8)
    text_w = text_right - left

    # 标题起点：取 OCR 标题框顶部，但给足够高度
    if outline.title_bbox_px:
        _, ttop, _, _ = m.rect(outline.title_bbox_px)
    else:
        ttop = 1.4
    title_lines = [s for s in outline.title.split(" ") if s] or [outline.title]
    n_lines = max(1, len(title_lines))
    title_size = 30 if n_lines >= 2 else 34
    title_h = n_lines * title_size * 1.25 / 72.0 + 0.2
    _multiline_textbox(slide, title_lines, left, ttop, text_w, title_h,
                       size=title_size, color=outline.text_color, bold=True,
                       anchor=MSO_ANCHOR.TOP)
    y = ttop + title_h + 0.12

    # 分隔线
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(left), Inches(y), Inches(left + 2.2), Inches(y))
    line.line.color.rgb = _hex(outline.accent_color, (39, 166, 139)); line.line.width = Pt(3)
    y += 0.22

    # 副标题
    if outline.subtitle:
        _textbox(slide, outline.subtitle, left, y, text_w, 0.5,
                 size=18, color=outline.accent_color, anchor=MSO_ANCHOR.TOP)
        y += 0.7

    # 信息块：竖向堆叠
    y = max(y, SLIDE_H * 0.55)
    for node in outline.nodes:
        txt = " ".join(x for x in [node.heading, node.subtext, node.body] if x).strip()
        if not txt:
            continue
        _textbox(slide, txt, left, y, text_w, 0.45,
                 size=13, color=outline.text_color, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.5


def _resolve_overlaps(placed, gap=0.06):
    """消除卡片矩形重叠：同行(纵向相近)按比例横向收缩并紧排，避免文字撞框。

    placed: [(order, node, [l,t,w,h]), ...]，原地修改 rect。
    """
    if len(placed) < 2:
        return
    # 按行分组（纵向中心相近）
    items = sorted(placed, key=lambda p: p[2][1])
    rows = []
    for it in items:
        cy = it[2][1] + it[2][3] / 2
        for row in rows:
            ref = row[0]
            rcy = ref[2][1] + ref[2][3] / 2
            if abs(cy - rcy) < max(ref[2][3], it[2][3]) * 0.6:
                row.append(it); break
        else:
            rows.append([it])
    for row in rows:
        if len(row) < 2:
            continue
        row.sort(key=lambda p: p[2][0])
        # 检测横向重叠，若有则在行内可用宽度上均分紧排
        overlap = False
        for a, b in zip(row, row[1:]):
            if a[2][0] + a[2][2] > b[2][0] + 0.02:
                overlap = True; break
        if not overlap:
            continue
        left = min(it[2][0] for it in row)
        right = max(it[2][0] + it[2][2] for it in row)
        n = len(row)
        total_gap = gap * (n - 1)
        cell_w = max(0.8, (right - left - total_gap) / n)
        x = left
        for it in row:
            it[2][0] = x
            it[2][2] = cell_w
            x += cell_w + gap


def _detect_table_nodes(outline):
    """识别表格节点，不依赖 VLM 的 role 命名（每次运行可能叫 table_header/
    table_row/table_section 等）。返回 (table_nodes, kind)：
      kind='multi'  —— 多个并列行节点（role 含 'table'/'row'）
      kind='flat'   —— 单个节点把整张表压成了 body_lines 列表
    无表格返回 (None, None)。
    """
    multi = [n for n in outline.nodes
             if "table" in n.role.lower() or n.role.lower().endswith("row")]
    real_rows = [n for n in multi if (n.heading or n.body)]
    if len(real_rows) >= 2:
        return real_rows, "multi"
    # 扁平表：某节点 role 含 table 且 body_lines 很多
    for n in outline.nodes:
        if "table" in n.role.lower() and len(n.body_lines) >= 4:
            return [n], "flat"
    return None, None


def _reconstruct_grid(flat_node):
    """从扁平 body_lines 重建网格。

    典型顺序：[col1..colN 表头][row1标签 row2标签 ...][row1末列 row2末列 ...]
    用「短词=表头/标签，长句=说明」的长度特征切分，输出二维 cells（含表头行）。
    失败返回 None。
    """
    toks = [bl.text.strip() for bl in flat_node.body_lines if bl.text.strip()]
    if len(toks) < 4:
        return None
    # 表头：前缀里的短词（<=6 字、不含句号）
    headers = []
    for t in toks:
        if len(t) <= 8 and "。" not in t and "，" not in t:
            headers.append(t)
        else:
            break
    ncol = len(headers)
    if ncol < 2:
        return None
    rest = toks[ncol:]
    # rest 里：含句号的是说明句(末列)，其余是行标签(首列)，按行序配对。
    descs = [t for t in rest if "。" in t]
    labels = [t for t in rest if "。" not in t]
    nrow = max(len(labels), len(descs))
    if nrow == 0:
        return None
    grid = [headers]
    for i in range(nrow):
        row = [""] * ncol
        if i < len(labels):
            row[0] = labels[i]           # 首列：阶段标签
        if i < len(descs) and ncol >= 2:
            row[ncol - 1] = descs[i]     # 末列：经济学意义说明
        grid.append(row)
    return grid


def _style_table_cell(cell, text, *, is_header, style, text_color):
    cell.text = text or ""
    cell.fill.solid()
    cell.fill.fore_color.rgb = _hex(
        style.header_band_color if is_header else style.card_fill,
        (39, 166, 139) if is_header else (28, 51, 70))
    if not is_header and style.card_style == "translucent":
        _set_alpha(cell.fill.fore_color, style.card_opacity)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(12.5 if is_header else 11.5)
            r.font.bold = is_header
            r.font.name = style.font_name
            r.font.color.rgb = _hex(text_color, (235, 235, 235))


def _render_table(slide, outline, m, style):
    """把表格节点渲染成对齐网格（role 名无关）。返回需继续走卡片的其余节点；
    无表格则返回 False。"""
    table_nodes, kind = _detect_table_nodes(outline)
    if not table_nodes:
        return False
    other = [n for n in outline.nodes if n not in table_nodes]

    xs0 = min(n.bbox_px[0] for n in table_nodes if n.bbox_px)
    ys0 = min(n.bbox_px[1] for n in table_nodes if n.bbox_px)
    xs1 = max(n.bbox_px[2] for n in table_nodes if n.bbox_px)
    ys1 = max(n.bbox_px[3] for n in table_nodes if n.bbox_px)
    l, t, w, h = m.rect([xs0, ys0, xs1, ys1])

    if kind == "flat":
        grid = _reconstruct_grid(table_nodes[0])
        if not grid:
            return False  # 重建失败，回退普通渲染
        nrow, ncol = len(grid), len(grid[0])
        gtbl = slide.shapes.add_table(nrow, ncol, Inches(l), Inches(t),
                                      Inches(w), Inches(h)).table
        for ri, row in enumerate(grid):
            for ci, val in enumerate(row):
                _style_table_cell(gtbl.cell(ri, ci), val,
                                  is_header=(ri == 0), style=style,
                                  text_color=outline.text_color)
        return other

    # multi：每个节点一行，两列（标签 / 说明）
    rows = sorted(table_nodes, key=lambda n: n.bbox_px[1] if n.bbox_px else 0)
    gtbl = slide.shapes.add_table(len(rows), 2, Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table
    gtbl.columns[0].width = Inches(w * 0.4)
    gtbl.columns[1].width = Inches(w * 0.6)
    for ri, node in enumerate(rows):
        is_header = "header" in node.role.lower() and not node.body
        _style_table_cell(gtbl.cell(ri, 0), node.heading, is_header=is_header,
                          style=style, text_color=outline.text_color)
        _style_table_cell(gtbl.cell(ri, 1), node.body, is_header=is_header,
                          style=style, text_color=outline.text_color)
    return other


def render_outline_to_slide(prs, outline) -> None:
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, outline.bg_color)
    _draw_grid(slide, outline.accent_color)
    m = _Mapper(outline.page_width, outline.page_height)
    style = outline.style

    # 封面专属布局
    lt = outline.layout_type.value if hasattr(outline.layout_type, "value") else str(outline.layout_type)
    if lt in ("title_cover", "image_content"):
        _render_cover(prs, slide, outline, m, style)
        return

    # 内容页：先放保留的复杂图形像素（底层），再叠加重生成的卡片/文字
    _place_keep_regions(slide, outline, m)

    # 标题 / 副标题（按真实坐标）
    if outline.title and outline.title_bbox_px:
        l, t, w, h = m.rect(outline.title_bbox_px)
        _textbox(slide, outline.title, l, t, w, h,
                 size=min(style.title_size, m.pt_from_h(outline.title_bbox_px[3] - outline.title_bbox_px[1])),
                 color=outline.text_color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if outline.subtitle and outline.subtitle_bbox_px:
        l, t, w, h = m.rect(outline.subtitle_bbox_px)
        _textbox(slide, outline.subtitle, l, t, w, h,
                 size=min(18, m.pt_from_h(outline.subtitle_bbox_px[3] - outline.subtitle_bbox_px[1])),
                 color=outline.accent_color, anchor=MSO_ANCHOR.MIDDLE)

    # 表格检测（role 名无关）：把表格节点渲染成对齐网格，其余节点继续走卡片
    nodes_for_cards = outline.nodes
    remaining = _render_table(slide, outline, m, style)
    if remaining is not False:
        nodes_for_cards = remaining

    # 节点卡片（按真实外接框，留 padding，应用自适应样式）
    pad_x, pad_y = 0.12, 0.1
    # 先把所有节点映射成矩形
    placed = []  # (order, node, [l,t,w,h])
    for node in nodes_for_cards:
        if not node.bbox_px:
            continue
        l, t, w, h = m.rect(node.bbox_px)
        l -= pad_x; t -= pad_y; w += 2 * pad_x; h += 2 * pad_y
        placed.append((node.order, node, [max(0, l), max(0, t), w, h]))

    # 行归一化：纵向相近的卡片归为一行，统一 top 与 height（消除大小悬殊）
    placed.sort(key=lambda p: p[2][1])
    rows = []
    for item in placed:
        t = item[2][1]
        h = item[2][3]
        for row in rows:
            rt = row[0][2][1]
            if abs(t - rt) < max(0.6, h * 0.5):
                row.append(item)
                break
        else:
            rows.append([item])
    for row in rows:
        if len(row) < 2:
            continue
        top = min(it[2][1] for it in row)
        bottom = max(it[2][1] + it[2][3] for it in row)
        for it in row:
            it[2][1] = top
            it[2][3] = bottom - top

    # 消除横向重叠（卡片溢出/相撞）
    _resolve_overlaps(placed)

    rects = {}
    for order, node, (l, t, w, h) in placed:
        card = _card(slide, l, t, w, h, style)
        if style.header_band and node.heading:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(l), Inches(t), Inches(w), Inches(0.06))
            band.fill.solid(); band.fill.fore_color.rgb = _hex(style.header_band_color, (39, 166, 139))
            band.line.fill.background(); band.shadow.inherit = False
        _fill_card_text(card, node, outline.text_color, outline.accent_color, style)
        rects[order] = (l, t, w, h)

    # 连接箭头
    for conn in outline.connections:
        if conn.from_order in rects and conn.to_order in rects:
            _connect(slide, rects[conn.from_order], rects[conn.to_order], style)

    # 底部说明条（按真实坐标，描边框）
    if outline.footer_note and outline.footer_bbox_px:
        l, t, w, h = m.rect(outline.footer_bbox_px)
        l -= pad_x; t -= pad_y; w += 2 * pad_x; h += 2 * pad_y
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(max(0, l)), Inches(max(0, t)), Inches(w), Inches(h))
        bar.fill.background()
        bar.line.color.rgb = _hex(outline.accent_color, (39, 166, 139)); bar.line.width = Pt(1.2)
        bar.shadow.inherit = False
        _set_radius(bar, 0.5)
        tf = bar.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = outline.footer_note
        r.font.size = Pt(11); r.font.color.rgb = _hex(outline.text_color); r.font.name = style.font_name
