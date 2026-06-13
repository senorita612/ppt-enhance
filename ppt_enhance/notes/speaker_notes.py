"""Generate speaker notes from SlideIR and write them into PPTX files."""

from __future__ import annotations

import json
import re
import zipfile
from collections import Counter
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Iterable, Sequence
from xml.etree import ElementTree

from pptx import Presentation

from ppt_enhance.agents.llm_client import LLMClient
from ppt_enhance.schemas.slide_ir import SlideIR, SlidePage


@dataclass
class SpeakerNote:
    page_no: int
    title: str
    note_text: str
    estimated_seconds: int
    reference_excerpt: str = ""
    generation_mode: str = "fallback"
    model: str = ""

    def to_dict(self) -> dict:
        return asdict(self)


@dataclass
class PageNarrative:
    page_no: int
    title: str
    intent: str
    key_message: str
    transition_from_previous: str = ""

    def to_dict(self) -> dict:
        return asdict(self)


@dataclass
class NarrativePlan:
    opening: str
    throughline: str
    closing: str
    pages: list[PageNarrative]
    generation_mode: str = "fallback"
    model: str = ""

    def to_dict(self) -> dict:
        data = asdict(self)
        data["pages"] = [page.to_dict() for page in self.pages]
        return data


@dataclass
class SpeakerNotesResult:
    notes: list[SpeakerNote]
    plan: NarrativePlan
    generation_mode: str
    model: str = ""

    def to_dict(self) -> dict:
        return {
            "generation_mode": self.generation_mode,
            "model": self.model,
            "plan": self.plan.to_dict(),
            "notes": [note.to_dict() for note in self.notes],
        }


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    try:
        import fitz
    except Exception:
        return ""

    doc = fitz.open(path)
    try:
        return "\n\n".join(page.get_text("text").strip() for page in doc if page.get_text("text").strip())
    finally:
        doc.close()


def _read_docx(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as zf:
            xml = zf.read("word/document.xml")
    except Exception:
        return ""

    root = ElementTree.fromstring(xml)
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: list[str] = []
    for para in root.findall(".//w:p", ns):
        parts = [node.text or "" for node in para.findall(".//w:t", ns)]
        text = "".join(parts).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def extract_reference_text(path: str | Path) -> str:
    """Extract plain text from a supplemental material file."""
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(p)
    if suffix == ".docx":
        return _read_docx(p)
    return _read_text_file(p)


def load_reference_text(
    paths: Iterable[str | Path] | None = None,
    raw_text: str | None = None,
) -> str:
    """Merge user-provided notes material into one plain-text corpus."""
    chunks: list[str] = []
    if raw_text and raw_text.strip():
        chunks.append(raw_text.strip())
    for path in paths or []:
        text = extract_reference_text(path).strip()
        if text:
            chunks.append(text)
    return "\n\n".join(chunks)


def _clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    text = re.sub(r"^[\-\*\u2022\u00b7\u25cf\u25aa\s]+", "", text)
    return text


def _page_texts(page: SlidePage) -> list[str]:
    texts = []
    for elem in page.text_elements():
        text = _clean_text(elem.final_text)
        if text:
            texts.append(text)
    return texts


def _page_title(page: SlidePage) -> str:
    titles = [
        _clean_text(elem.final_text)
        for elem in page.text_elements()
        if getattr(elem.type, "value", elem.type) == "title" and elem.final_text.strip()
    ]
    if titles:
        return titles[0]
    texts = _page_texts(page)
    return texts[0] if texts else f"第 {page.page_no} 页"


def _split_reference(text: str, chunk_size: int = 900) -> list[str]:
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    if not paragraphs:
        return []

    chunks: list[str] = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) + 2 <= chunk_size:
            current = f"{current}\n\n{para}".strip()
        else:
            if current:
                chunks.append(current)
            if len(para) <= chunk_size:
                current = para
            else:
                for i in range(0, len(para), chunk_size):
                    chunks.append(para[i : i + chunk_size])
                current = ""
    if current:
        chunks.append(current)
    return chunks


def _terms(text: str) -> Counter:
    words = re.findall(r"[A-Za-z][A-Za-z0-9_\-]{1,}", text.lower())
    cjk_sequences = re.findall(r"[\u4e00-\u9fff]{2,}", text)
    terms: list[str] = words
    for seq in cjk_sequences:
        terms.extend(seq[i : i + 2] for i in range(max(0, len(seq) - 1)))
        terms.extend(seq[i : i + 3] for i in range(max(0, len(seq) - 2)))
    return Counter(terms)


def _retrieve_context(query: str, reference_text: str, limit: int = 2) -> str:
    chunks = _split_reference(reference_text)
    if not chunks:
        return ""

    query_terms = _terms(query)
    if not query_terms:
        return "\n\n".join(chunks[:limit])

    scored: list[tuple[float, str]] = []
    for chunk in chunks:
        chunk_terms = _terms(chunk)
        score = sum(min(count, chunk_terms.get(term, 0)) for term, count in query_terms.items())
        if score:
            scored.append((float(score), chunk))

    if not scored:
        return "\n\n".join(chunks[:limit])

    scored.sort(key=lambda item: item[0], reverse=True)
    return "\n\n".join(chunk for _, chunk in scored[:limit])


def _first_sentence(text: str, max_len: int = 120) -> str:
    text = _clean_text(text)
    if not text:
        return ""
    parts = re.split(r"(?<=[。！？.!?])", text)
    sentence = parts[0].strip() if parts else text
    return sentence[:max_len].strip()


def _deck_summary(slide_ir: SlideIR) -> list[dict]:
    summary = []
    for page in sorted(slide_ir.pages, key=lambda p: p.page_no):
        texts = _page_texts(page)
        summary.append(
            {
                "page_no": page.page_no,
                "title": _page_title(page),
                "texts": texts[:8],
            }
        )
    return summary


def _fallback_plan(slide_ir: SlideIR, style: str) -> NarrativePlan:
    pages = sorted(slide_ir.pages, key=lambda p: p.page_no)
    titles = [_page_title(page) for page in pages]
    throughline = " -> ".join(titles[:6]) if titles else "逐页讲解文稿内容"
    plan_pages: list[PageNarrative] = []

    for index, page in enumerate(pages):
        title = _page_title(page)
        texts = [text for text in _page_texts(page) if text != title]
        key_message = texts[0] if texts else title
        if index == 0:
            transition = "先交代本次演示的主题和听众需要关注的问题。"
        else:
            transition = f"上一页讲到 {titles[index - 1]}，这一页继续展开到 {title}。"
        plan_pages.append(
            PageNarrative(
                page_no=page.page_no,
                title=title,
                intent=f"用{style}的口吻解释本页核心内容，并让听众理解它在整份材料中的作用。",
                key_message=key_message,
                transition_from_previous=transition,
            )
        )

    opening = f"本次演示围绕 {titles[0]} 展开。" if titles else "本次演示围绕文稿主题展开。"
    closing = "最后回到整份材料的核心结论，提醒听众下一步可以关注或行动的方向。"
    return NarrativePlan(
        opening=opening,
        throughline=throughline,
        closing=closing,
        pages=plan_pages,
    )


def _coerce_page_narrative(data: dict, fallback: PageNarrative) -> PageNarrative:
    try:
        page_no = int(data.get("page_no", fallback.page_no))
    except (TypeError, ValueError):
        page_no = fallback.page_no
    return PageNarrative(
        page_no=page_no,
        title=str(data.get("title") or fallback.title),
        intent=str(data.get("intent") or fallback.intent),
        key_message=str(data.get("key_message") or fallback.key_message),
        transition_from_previous=str(
            data.get("transition_from_previous") or fallback.transition_from_previous
        ),
    )


def _llm_plan(
    slide_ir: SlideIR,
    reference_text: str,
    seconds_per_slide: int,
    style: str,
    llm: LLMClient,
) -> NarrativePlan | None:
    system = (
        "你是演示文稿总撰稿人。请先为整份 PPT 设计演讲叙事计划，再供逐页讲稿使用。"
        "计划要体现整场演讲的主线、每页讲述目的、核心信息和页间过渡。"
        "只输出 JSON。"
    )
    reference_preview = reference_text[:5000]
    payload = {
        "style": style,
        "seconds_per_slide": seconds_per_slide,
        "slides": _deck_summary(slide_ir),
        "protected_terms": slide_ir.protected_terms,
        "reference_material": reference_preview,
        "json_schema": {
            "opening": "整场演讲开场思路",
            "throughline": "整份 PPT 的叙事主线",
            "closing": "整场演讲收束方式",
            "pages": [
                {
                    "page_no": 1,
                    "title": "页标题",
                    "intent": "这一页在整场演讲中的任务",
                    "key_message": "这一页必须讲清楚的一句话",
                    "transition_from_previous": "从上一页过渡到本页的话",
                }
            ],
        },
    }
    try:
        result = llm.chat_json(system, json.dumps(payload, ensure_ascii=False), temperature=0.25)
    except Exception:
        return None

    fallback = _fallback_plan(slide_ir, style)
    fallback_by_page = {page.page_no: page for page in fallback.pages}
    pages: list[PageNarrative] = []
    for item in result.get("pages", []) or []:
        try:
            page_no = int(item.get("page_no", 0))
        except (TypeError, ValueError):
            continue
        fallback_page = fallback_by_page.get(page_no)
        if fallback_page:
            pages.append(_coerce_page_narrative(item, fallback_page))
    if not pages:
        return None

    covered = {page.page_no for page in pages}
    for fallback_page in fallback.pages:
        if fallback_page.page_no not in covered:
            pages.append(fallback_page)
    pages.sort(key=lambda page: page.page_no)

    return NarrativePlan(
        opening=str(result.get("opening") or fallback.opening),
        throughline=str(result.get("throughline") or fallback.throughline),
        closing=str(result.get("closing") or fallback.closing),
        pages=pages,
        generation_mode="llm",
        model=llm.model,
    )


def _fallback_note(
    page: SlidePage,
    reference_excerpt: str,
    seconds_per_slide: int,
    narrative: PageNarrative,
) -> str:
    title = _page_title(page)
    points = [p for p in _page_texts(page) if p != title]
    target_points = max(2, min(5, seconds_per_slide // 18))

    lines = [
        narrative.transition_from_previous,
        f"这一页我们进入 {title}。",
        f"这一页的重点是：{narrative.key_message}。",
        f"这里要完成的讲述任务是：{narrative.intent}",
    ]
    if points:
        first = points[0]
        lines.append(f"可以先把页面上的第一层信息讲清楚：{first}。")
        for point in points[1:target_points]:
            lines.append(f"接着展开 {point}，说明它如何支撑本页结论。")
    else:
        lines.append("这里可以先概括页面主题，再解释为什么这一页值得单独展开。")

    extra = _first_sentence(reference_excerpt)
    if extra:
        lines.append(f"结合补充资料，可以补充一句：{extra}")

    lines.append("最后用一句话把本页结论落到听众能带走的信息上，再进入下一页。")
    return "\n".join(line for line in lines if line.strip())


def _llm_note(
    page: SlidePage,
    reference_excerpt: str,
    seconds_per_slide: int,
    style: str,
    llm: LLMClient,
    protected_terms: Sequence[str],
    narrative: PageNarrative,
    plan: NarrativePlan,
) -> str | None:
    system = (
        "你是中文演示文稿讲稿撰写助手。请根据每页 PPT 内容生成 Speaker Notes，"
        "讲稿要口语自然、信息准确、便于演讲者直接照着讲。不得改写专有名词、数字、公式。"
        "只输出 JSON：{\"note\":\"...\"}。"
    )
    payload = {
        "page_no": page.page_no,
        "title": _page_title(page),
        "slide_texts": _page_texts(page),
        "reference_excerpt": reference_excerpt,
        "seconds_per_slide": seconds_per_slide,
        "style": style,
        "deck_opening": plan.opening,
        "deck_throughline": plan.throughline,
        "deck_closing": plan.closing,
        "page_intent": narrative.intent,
        "page_key_message": narrative.key_message,
        "transition_from_previous": narrative.transition_from_previous,
        "protected_terms": list(protected_terms),
        "requirements": [
            "写成演讲者可以直接照读的讲稿，而不是摘要或 bullet point",
            "开头自然承接上一页，结尾为下一页留出过渡",
            "内容要围绕 page_intent 和 page_key_message 展开",
            "可以适度补充解释，但不得编造 slide_texts 和 reference_excerpt 之外的事实",
            "长度匹配目标时长",
        ],
    }
    try:
        result = llm.chat_json(system, json.dumps(payload, ensure_ascii=False), temperature=0.35)
    except Exception:
        return None
    note = str(result.get("note", "")).strip()
    return note or None


def generate_speaker_notes(
    slide_ir: SlideIR,
    reference_text: str = "",
    seconds_per_slide: int = 75,
    style: str = "课程讲解",
    llm: LLMClient | None = None,
) -> SpeakerNotesResult:
    """Generate a deck-level narrative plan and one speaker note per slide."""
    llm = llm or LLMClient()
    plan = None
    if llm.available:
        plan = _llm_plan(slide_ir, reference_text, seconds_per_slide, style, llm)
    if plan is None:
        plan = _fallback_plan(slide_ir, style)

    narrative_by_page = {page.page_no: page for page in plan.pages}
    notes: list[SpeakerNote] = []
    llm_note_count = 0

    for page in sorted(slide_ir.pages, key=lambda p: p.page_no):
        title = _page_title(page)
        query = "\n".join([title, *_page_texts(page), *slide_ir.protected_terms])
        reference_excerpt = _retrieve_context(query, reference_text)
        narrative = narrative_by_page.get(
            page.page_no,
            PageNarrative(
                page_no=page.page_no,
                title=title,
                intent=f"解释 {title} 的核心内容。",
                key_message=title,
            ),
        )

        note_text = None
        if llm.available:
            note_text = _llm_note(
                page,
                reference_excerpt,
                seconds_per_slide,
                style,
                llm,
                slide_ir.protected_terms,
                narrative,
                plan,
            )
        note_mode = "llm" if note_text else "fallback"
        if not note_text:
            note_text = _fallback_note(page, reference_excerpt, seconds_per_slide, narrative)
        else:
            llm_note_count += 1

        notes.append(
            SpeakerNote(
                page_no=page.page_no,
                title=title,
                note_text=note_text,
                estimated_seconds=seconds_per_slide,
                reference_excerpt=reference_excerpt,
                generation_mode=note_mode,
                model=llm.model if note_mode == "llm" else "",
            )
        )

    if llm_note_count == len(notes) and plan.generation_mode == "llm":
        mode = "llm"
    elif llm_note_count:
        mode = "mixed"
    else:
        mode = "fallback"
    model = llm.model if llm_note_count or plan.generation_mode == "llm" else ""
    return SpeakerNotesResult(notes=notes, plan=plan, generation_mode=mode, model=model)


def write_speaker_notes(
    pptx_path: str | Path,
    notes: Sequence[SpeakerNote],
    output_path: str | Path | None = None,
) -> Path:
    """Write generated speaker notes into a PPTX file."""
    pptx_path = Path(pptx_path)
    output_path = Path(output_path) if output_path else pptx_path
    note_by_page = {note.page_no: note.note_text for note in notes}

    prs = Presentation(str(pptx_path))
    for index, slide in enumerate(prs.slides, start=1):
        text = note_by_page.get(index, "").strip()
        if not text:
            continue
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        notes_frame.text = text

    prs.save(str(output_path))
    return output_path


def save_speaker_notes_json(notes: Sequence[SpeakerNote] | SpeakerNotesResult, path: str | Path) -> Path:
    path = Path(path)
    data = notes.to_dict() if isinstance(notes, SpeakerNotesResult) else [note.to_dict() for note in notes]
    path.write_text(
        json.dumps(data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return path


def speaker_notes_to_markdown(notes: Sequence[SpeakerNote] | SpeakerNotesResult) -> str:
    plan = notes.plan if isinstance(notes, SpeakerNotesResult) else None
    note_items = notes.notes if isinstance(notes, SpeakerNotesResult) else notes
    sections = []
    if plan:
        sections.append(
            "# Speaker Notes\n\n"
            f"- Generation mode: {notes.generation_mode}\n"
            f"- Model: {notes.model or 'local fallback'}\n"
            f"- Opening: {plan.opening}\n"
            f"- Throughline: {plan.throughline}\n"
            f"- Closing: {plan.closing}"
        )
    for note in note_items:
        title = note.title or f"Slide {note.page_no}"
        sections.append(
            f"## Page {note.page_no}: {title}\n\n"
            f"Mode: {note.generation_mode}"
            f"{f' ({note.model})' if note.model else ''}\n\n"
            f"{note.note_text.strip()}"
        )
    return "\n\n".join(sections) + ("\n" if sections else "")


def save_speaker_notes_markdown(notes: Sequence[SpeakerNote] | SpeakerNotesResult, path: str | Path) -> Path:
    path = Path(path)
    path.write_text(speaker_notes_to_markdown(notes), encoding="utf-8")
    return path
