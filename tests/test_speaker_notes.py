from pptx import Presentation

from ppt_enhance.notes import SpeakerNote, generate_speaker_notes, write_speaker_notes
from ppt_enhance.schemas.slide_ir import BBox, ElementType, SlideElement, SlideIR, SlidePage


class OfflineLLM:
    available = False


class FakeLLM:
    available = True
    model = "test-model"

    def __init__(self):
        self.calls = []

    def chat_json(self, system: str, user: str, temperature: float = 0.2):
        import json

        self.calls.append(system)
        payload = json.loads(user)
        if "slides" in payload:
            return {
                "opening": "先说明 Transformer 为什么重要。",
                "throughline": "从结构到机制再到应用。",
                "closing": "最后回到建模能力的提升。",
                "pages": [
                    {
                        "page_no": 1,
                        "title": "Transformer 架构",
                        "intent": "解释这一页在整场演讲中的定位。",
                        "key_message": "Transformer 的核心是自注意力机制。",
                        "transition_from_previous": "开场后我们先看整体架构。",
                    }
                ],
            }
        return {"note": "开场后我们先看整体架构。这一页要讲清楚 Transformer 的核心是自注意力机制。"}


def _sample_ir() -> SlideIR:
    title = SlideElement(
        id="title_1",
        type=ElementType.TITLE,
        text="Transformer 架构",
        bbox=BBox(x0=10, y0=10, x1=400, y1=70),
        page_no=1,
    )
    body = SlideElement(
        id="body_1",
        type=ElementType.TEXT,
        text="自注意力机制用于建模序列内部关系",
        bbox=BBox(x0=10, y0=100, x1=700, y1=140),
        page_no=1,
    )
    return SlideIR(
        source_pdf="sample.pdf",
        pages=[SlidePage(page_no=1, width=960, height=540, elements=[title, body])],
        protected_terms=["Transformer"],
    )


def test_generate_speaker_notes_fallback():
    result = generate_speaker_notes(
        _sample_ir(),
        reference_text="补充资料：Transformer 在 2017 年提出，核心是 Self-Attention。",
        seconds_per_slide=60,
        llm=OfflineLLM(),
    )

    notes = result.notes
    assert result.generation_mode == "fallback"
    assert len(notes) == 1
    assert notes[0].page_no == 1
    assert "Transformer 架构" in notes[0].note_text
    assert "自注意力机制" in notes[0].note_text


def test_generate_speaker_notes_uses_deck_plan():
    fake = FakeLLM()
    result = generate_speaker_notes(
        _sample_ir(),
        reference_text="补充资料：Transformer 在 2017 年提出，核心是 Self-Attention。",
        seconds_per_slide=60,
        llm=fake,
    )

    assert result.generation_mode == "llm"
    assert result.model == "test-model"
    assert result.plan.throughline == "从结构到机制再到应用。"
    assert result.notes[0].generation_mode == "llm"
    assert "自注意力机制" in result.notes[0].note_text
    assert len(fake.calls) == 2


def test_write_speaker_notes(tmp_path):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    prs.slides.add_slide(blank)
    prs.slides.add_slide(blank)

    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    write_speaker_notes(
        pptx_path,
        [
            SpeakerNote(page_no=1, title="第一页", note_text="第一页讲稿", estimated_seconds=60),
            SpeakerNote(page_no=2, title="第二页", note_text="第二页讲稿", estimated_seconds=60),
        ],
    )

    reopened = Presentation(str(pptx_path))
    assert "第一页讲稿" in reopened.slides[0].notes_slide.notes_text_frame.text
    assert "第二页讲稿" in reopened.slides[1].notes_slide.notes_text_frame.text
