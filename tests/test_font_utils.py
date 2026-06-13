from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_enhance.builder.font_utils import (
    default_cjk_font,
    set_paragraph_cjk_font,
    set_run_cjk_font,
)


def _typeface(run, slot: str) -> str | None:
    r_pr = run._r.rPr
    node = r_pr.find(qn(f"a:{slot}")) if r_pr is not None else None
    return node.get("typeface") if node is not None else None


def test_set_run_cjk_font_writes_east_asian_typeface():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = "\u4e2d\u6587 Test"

    resolved = set_run_cjk_font(run, "PingFang SC")

    assert resolved == default_cjk_font()
    assert run.font.name == default_cjk_font()
    assert _typeface(run, "latin") == default_cjk_font()
    assert _typeface(run, "ea") == default_cjk_font()
    assert _typeface(run, "cs") == default_cjk_font()


def test_set_paragraph_cjk_font_updates_existing_runs():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    paragraph = textbox.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "\u6f14\u793a"

    resolved = set_paragraph_cjk_font(paragraph)

    assert resolved == default_cjk_font()
    assert paragraph.font.name == default_cjk_font()
    assert _typeface(run, "ea") == default_cjk_font()
